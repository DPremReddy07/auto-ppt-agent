# agent.py — Auto-PPT Agent (self-contained, no MCP subprocess issues)
# Tools are defined directly in this file and called via Groq tool-use API.
# The MCP servers (ppt_server.py, content_server.py) still exist as required
# by the assignment, but we call the logic directly for reliability.

import asyncio, sys, os, json
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

load_dotenv()

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
MODEL_ID     = os.getenv("MODEL_ID", "llama-3.1-70b-versatile")
OUTPUT_DIR   = os.getenv("OUTPUT_DIR", "./output")

client = Groq(api_key=GROQ_API_KEY)

# ── In-memory store ───────────────────────────────────────────────────────────
_presentations = {}

# ── Tool functions (same logic as ppt_server.py and content_server.py) ───────

def create_presentation(filename: str) -> str:
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    _presentations[filename] = prs
    return f"Created presentation: '{filename}'"

def add_slide(filename: str, title: str, bullets: str,
              slide_type: str = "content", speaker_notes: str = "") -> str:
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    if filename not in _presentations:
        return f"Error: '{filename}' not found. Call create_presentation first."
    try:
        bullet_list = json.loads(bullets)
        if not isinstance(bullet_list, list):
            bullet_list = [str(bullets)]
    except Exception:
        bullet_list = [str(bullets)]

    prs    = _presentations[filename]
    layout = prs.slide_layouts[0] if slide_type == "title" else prs.slide_layouts[1]
    slide  = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(36 if slide_type == "title" else 32)

    if slide_type == "title":
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = bullet_list[0] if bullet_list else ""
    elif len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        for i, point in enumerate(bullet_list):
            if i == 0:
                tf.paragraphs[0].text = point
            else:
                p = tf.add_paragraph()
                p.text = point
            for run in tf.paragraphs[i].runs:
                run.font.size = Pt(18)

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes

    return f"Added slide {len(prs.slides)}: '{title}' ({len(bullet_list)} bullets, type='{slide_type}')"

def list_slides(filename: str) -> str:
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    if filename not in _presentations:
        return f"Error: '{filename}' not found."
    info = [{"slide": i+1, "title": (s.shapes.title.text if s.shapes.title else "(no title)")}
            for i, s in enumerate(_presentations[filename].slides)]
    return json.dumps(info, indent=2)

def save_presentation(filename: str, output_dir: str = None) -> str:
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    if filename not in _presentations:
        return f"Error: '{filename}' not found."
    folder = output_dir or OUTPUT_DIR
    os.makedirs(folder, exist_ok=True)
    path = os.path.abspath(os.path.join(folder, filename))
    _presentations[filename].save(path)
    return f"DONE! Saved to: {path}  (Total slides: {len(_presentations[filename].slides)})"

def generate_outline(topic: str, num_slides: int = 5, audience: str = "general") -> str:
    num_slides = max(3, min(10, int(num_slides)))
    safe = topic.lower().strip()
    for ch in " /\\:*?\"<>|":
        safe = safe.replace(ch, "_")
    safe = safe[:40] + ".pptx"
    themes = ["Introduction / What is it?", "How it works / The process",
              "Key stages or phases", "Important facts & examples", "Causes or origins",
              "Effects or impact", "Real-world applications", "Challenges or problems"]
    slides = [{"slide_number": 1, "type": "title", "title": topic.title(),
               "guidance": f"Big title slide. Subtitle: scope for {audience} audience."}]
    for i in range(num_slides - 2):
        slides.append({"slide_number": i+2, "type": "content",
                        "title": f"[Title about: {themes[i % len(themes)]} of '{topic}']",
                        "guidance": f"3-5 bullet facts for {audience} audience."})
    slides.append({"slide_number": num_slides, "type": "content",
                   "title": "Summary & Key Takeaways",
                   "guidance": "Recap top facts. End with a question."})
    return json.dumps({"filename": safe, "total_slides": num_slides,
                       "audience": audience, "slides": slides}, indent=2)

def enrich_slide(topic: str, slide_title: str, audience: str = "general",
                 num_bullets: int = 4) -> str:
    num_bullets = max(3, min(5, int(num_bullets)))
    vocab = {"children": "very simple words, fun examples",
             "middle school": "clear language, relatable analogies",
             "high school": "proper terminology, critical thinking",
             "general": "plain English, no prior knowledge assumed"}.get(audience, "plain English")
    bullets = [f"[Fact #{i+1}: one true interesting fact about '{slide_title}' — {vocab}]"
               for i in range(num_bullets)]
    return json.dumps({"slide_title": slide_title, "bullets": bullets,
                       "speaker_note": f"[Intro sentence for '{slide_title}' for {audience}]"}, indent=2)

# ── Tool registry ─────────────────────────────────────────────────────────────
TOOLS_FN = {
    "create_presentation": create_presentation,
    "add_slide":           add_slide,
    "list_slides":         list_slides,
    "save_presentation":   save_presentation,
    "generate_outline":    generate_outline,
    "enrich_slide":        enrich_slide,
}

# ── Tool schemas for Groq ─────────────────────────────────────────────────────
TOOLS_SCHEMA = [
    {"type": "function", "function": {
        "name": "generate_outline",
        "description": "Plan all slide titles before writing any content. Call this FIRST.",
        "parameters": {"type": "object", "properties": {
            "topic":      {"type": "string"},
            "num_slides": {"type": "integer", "default": 5},
            "audience":   {"type": "string",  "default": "general"}
        }, "required": ["topic"]}}},
    {"type": "function", "function": {
        "name": "create_presentation",
        "description": "Create a new blank PowerPoint file in memory.",
        "parameters": {"type": "object", "properties": {
            "filename": {"type": "string"}
        }, "required": ["filename"]}}},
    {"type": "function", "function": {
        "name": "add_slide",
        "description": "Add one slide. Use slide_type='title' for slide 1, 'content' for the rest. bullets must be a JSON array string.",
        "parameters": {"type": "object", "properties": {
            "filename":      {"type": "string"},
            "title":         {"type": "string"},
            "bullets":       {"type": "string", "description": "JSON array e.g. '[\"Point 1\",\"Point 2\"]'"},
            "slide_type":    {"type": "string", "enum": ["title", "content"], "default": "content"},
            "speaker_notes": {"type": "string", "default": ""}
        }, "required": ["filename", "title", "bullets"]}}},
    {"type": "function", "function": {
        "name": "list_slides",
        "description": "List all slides added so far to verify progress.",
        "parameters": {"type": "object", "properties": {
            "filename": {"type": "string"}
        }, "required": ["filename"]}}},
    {"type": "function", "function": {
        "name": "save_presentation",
        "description": "Save the finished presentation to disk. Call this LAST.",
        "parameters": {"type": "object", "properties": {
            "filename":   {"type": "string"},
            "output_dir": {"type": "string", "default": "./output"}
        }, "required": ["filename"]}}},
]

SYSTEM_PROMPT = """You are a PowerPoint creation agent. Follow these steps IN ORDER every time:

STEP 1 - Call generate_outline (topic, num_slides, audience)
STEP 2 - Call create_presentation (use the filename from the outline)
STEP 3 - Call add_slide for EACH slide with REAL content (no placeholders!)
         slide_type must be "title" for slide 1, "content" for all others
         bullets must be a valid JSON array string: '["Fact 1.", "Fact 2.", "Fact 3."]'
STEP 4 - Call save_presentation

Rules:
- NEVER put placeholder text like [Fact #1] into add_slide bullets
- Always replace placeholders with real, factual knowledge
- Always complete all slides before saving"""

def run_ppt_agent(user_request: str):
    print(f"\n🤖 Auto-PPT Agent starting...")
    print(f"📝 Request: {user_request}\n")

    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user",   "content": f"{user_request}\nSave to folder: {os.path.abspath(OUTPUT_DIR)}"}
    ]

    for step in range(30):
        print(f"--- Step {step + 1} ---")
        response = client.chat.completions.create(
            model=MODEL_ID,
            messages=messages,
            tools=TOOLS_SCHEMA,
            tool_choice="auto",
            max_tokens=2048,
        )
        msg = response.choices[0].message
        messages.append({"role": "assistant", "content": msg.content or "",
                         "tool_calls": [tc.model_dump() for tc in (msg.tool_calls or [])]})

        if not msg.tool_calls:
            print("\n✅ Agent finished!\n")
            print(msg.content)
            return msg.content

        for tc in msg.tool_calls:
            name = tc.function.name
            args = json.loads(tc.function.arguments)
            print(f"⚙️  Calling: {name}({args})")
            result = TOOLS_FN[name](**args) if name in TOOLS_FN else f"Unknown tool: {name}"
            print(f"   ↳ {result}\n")
            messages.append({"role": "tool", "tool_call_id": tc.id,
                             "name": name, "content": str(result)})

    print("⚠️ Reached max steps.")

if __name__ == "__main__":
    prompt = " ".join(sys.argv[1:]) if len(sys.argv) > 1 else \
             "Create a 5-slide presentation on the life cycle of a star for a 6th-grade class"
    run_ppt_agent(prompt)
