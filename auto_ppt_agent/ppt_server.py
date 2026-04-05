# =============================================
# ppt_server.py  –  MCP Server #1: PowerPoint
# =============================================
# This file IS the MCP server.
# It exposes 4 tools the agent can call:
#   1. create_presentation  – start a new .pptx
#   2. add_slide            – add a slide with title + bullets
#   3. list_slides          – see what slides exist so far
#   4. save_presentation    – write the file to disk
#
# Run with:  python ppt_server.py
# (You don't run this manually – the agent starts it automatically)
# =============================================

import os
import json
from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── In-memory store (lives while server is running) ──────────────────────────
# Key  = filename (e.g. "stars.pptx")
# Value = python-pptx Presentation object
_presentations: dict = {}

# ── Create the FastMCP server (same pattern as your system_server.py) ────────
mcp = FastMCP("ppt-server")


# ─────────────────────────────────────────────────────────────────────────────
# TOOL 1 – Create a blank presentation
# ─────────────────────────────────────────────────────────────────────────────
@mcp.tool()
def create_presentation(filename: str) -> str:
    """
    Initialises a brand-new PowerPoint file and holds it in memory.
    Call this FIRST before adding any slides.

    Args:
        filename: The name to save the file as, e.g. 'life_of_a_star.pptx'
                  (the .pptx extension is added automatically if missing)

    Returns:
        A confirmation string.
    """
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    prs = Presentation()

    # Set slide size to widescreen (16:9) – more modern look
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _presentations[filename] = prs
    return f"✅ Created presentation: '{filename}' (16:9 widescreen). Ready to add slides."


# ─────────────────────────────────────────────────────────────────────────────
# TOOL 2 – Add a slide
# ─────────────────────────────────────────────────────────────────────────────
@mcp.tool()
def add_slide(
    filename: str,
    title: str,
    bullets: str,          # JSON string e.g. '["Point 1","Point 2","Point 3"]'
    slide_type: str = "content",   # "title" | "content" | "section"
    speaker_notes: str = ""
) -> str:
    """
    Adds one slide to the presentation.
    Call this once per slide, in order.

    Args:
        filename:      The name you used in create_presentation.
        title:         The slide heading (keep under 60 chars).
        bullets:       A JSON array string of 3-5 bullet points.
                       Example: '["Stars begin as nebulae","Gravity pulls gas together"]'
        slide_type:    "title"   → big centred title slide (use for slide 1 only)
                       "content" → title + bullet list (use for most slides)
                       "section" → bold section-break slide
        speaker_notes: Optional notes for the presenter.

    Returns:
        Confirmation or error message.
    """
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    if filename not in _presentations:
        return (
            f"❌ Presentation '{filename}' not found in memory. "
            "Did you call create_presentation first?"
        )

    # Parse bullets safely
    try:
        bullet_list = json.loads(bullets)
        if not isinstance(bullet_list, list):
            bullet_list = [str(bullets)]
    except (json.JSONDecodeError, TypeError):
        bullet_list = [str(bullets)]

    prs = _presentations[filename]

    # ── Pick the right slide layout ───────────────────────────────────────────
    # Layout index reference (standard python-pptx blank theme):
    #   0  = Title Slide  (big title + subtitle)
    #   1  = Title and Content (title + body placeholder)
    #   2  = Title Only
    #   5  = Blank
    if slide_type == "title":
        layout = prs.slide_layouts[0]    # Title slide
    elif slide_type == "section":
        layout = prs.slide_layouts[2]    # Title only → we'll add a big text box
    else:
        layout = prs.slide_layouts[1]    # Title + Content (default)

    slide = prs.slides.add_slide(layout)

    # ── Fill in the TITLE placeholder ─────────────────────────────────────────
    if slide.shapes.title:
        slide.shapes.title.text = title
        # Style the title font
        tf = slide.shapes.title.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(36) if slide_type == "title" else Pt(32)

    # ── Fill in the CONTENT placeholder (bullets) ─────────────────────────────
    if slide_type == "title":
        # For the title slide, put a single subtitle line
        subtitle = bullet_list[0] if bullet_list else ""
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    elif slide_type == "content" and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_list):
            if i == 0:
                tf.paragraphs[0].text = point
                tf.paragraphs[0].level = 0
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

            # Style each bullet
            current_para = tf.paragraphs[i]
            for run in current_para.runs:
                run.font.size = Pt(18)

    elif slide_type == "section":
        # Add a big centred text box for section break slides
        from pptx.util import Inches
        left   = Inches(1)
        top    = Inches(2.5)
        width  = Inches(11.33)
        height = Inches(2.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text      = "\n".join(bullet_list)
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(22)

    # ── Add speaker notes ──────────────────────────────────────────────────────
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes

    slide_number = len(prs.slides)
    return (
        f"✅ Added slide {slide_number}: '{title}' "
        f"({len(bullet_list)} points, type='{slide_type}')"
    )


# ─────────────────────────────────────────────────────────────────────────────
# TOOL 3 – List slides (so agent can check its own progress)
# ─────────────────────────────────────────────────────────────────────────────
@mcp.tool()
def list_slides(filename: str) -> str:
    """
    Returns a summary of all slides added so far.
    Useful for the agent to verify its own progress before saving.

    Args:
        filename: The presentation filename.

    Returns:
        JSON list of slide titles, or an error.
    """
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    if filename not in _presentations:
        return f"❌ Presentation '{filename}' not found."

    prs = _presentations[filename]
    slides_info = []
    for i, slide in enumerate(prs.slides, start=1):
        title_text = "(no title)"
        if slide.shapes.title and slide.shapes.title.text:
            title_text = slide.shapes.title.text
        slides_info.append({"slide": i, "title": title_text})

    return json.dumps(slides_info, indent=2)


# ─────────────────────────────────────────────────────────────────────────────
# TOOL 4 – Save presentation to disk
# ─────────────────────────────────────────────────────────────────────────────
@mcp.tool()
def save_presentation(filename: str, output_dir: str = "./output") -> str:
    """
    Saves the finished presentation to disk as a .pptx file.
    Call this LAST, after all slides have been added.

    Args:
        filename:   The name used in create_presentation.
        output_dir: Folder to save in (created if it doesn't exist).

    Returns:
        The full file path where the file was saved.
    """
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    if filename not in _presentations:
        return (
            f"❌ Presentation '{filename}' not found in memory. "
            "Cannot save a presentation that was never created."
        )

    os.makedirs(output_dir, exist_ok=True)
    full_path = os.path.join(output_dir, filename)

    try:
        _presentations[filename].save(full_path)
        abs_path = os.path.abspath(full_path)
        slide_count = len(_presentations[filename].slides)
        return (
            f"✅ DONE! Saved '{filename}' → {abs_path}\n"
            f"   Total slides: {slide_count}\n"
            f"   Open the file in PowerPoint or Google Slides to view it."
        )
    except Exception as e:
        return f"❌ Save failed: {e}"


# ─────────────────────────────────────────────────────────────────────────────
# Entry point
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    # This starts the stdio MCP server — the agent connects to it automatically
    print("🟢 PPT MCP Server is running…", flush=True)
    mcp.run()